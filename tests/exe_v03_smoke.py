"""Frozen-EXE runtime smoke for v0.3 (AI-enhanced conversion).

Built with the SAME collection flags as md-desk.exe (notably
``--collect-submodules markitdown_ocr`` and ``--hidden-import openai``), so it
exercises the identical bundled module set. It proves, INSIDE the frozen
binary (no venv):

  1. ``openai`` and the vendored ``markitdown_ocr`` are collected (importable).
  2. ``is_ocr_plugin_available()`` is True in the frozen build.
  3. AI-OFF path is byte-for-byte equivalent to plain ``MarkItDown()`` (v0.2
     behavior preserved in the frozen binary).
  4. The frozen plugin path (explicit ``import markitdown_ocr`` +
     ``register_converters`` because PyInstaller drops ``.dist-info`` entry
     points) registers the 4 OCR converters.
  5. A built-in LLM image description works end-to-end using a NO-NETWORK
     dummy OpenAI-compatible client.
  6. Scanned-PDF OCR works end-to-end via the same dummy client, with the
     OCR text surfaced as ``*[Image OCR]...[End OCR]*``.

A dummy client (no network) is used so the smoke is hermetic and fast.

stdout is forced to UTF-8. All printed diagnostics are ASCII.

Run: build the smoke EXE (build_v03_smoke.sh) then execute it.
"""

import io
import sys
import tempfile
from pathlib import Path

try:
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")
except Exception:
    pass

sys.path.insert(0, ".")

from markitdown import MarkItDown  # noqa: E402
import openai  # noqa: E402  (collection check)

from src.engine_config import EngineConfig, is_ocr_plugin_available  # noqa: E402
from src.markitdown_factory import MarkItDownFactory  # noqa: E402
from src.converter import convert_file  # noqa: E402


def _check(label, ok, detail=""):
    print(("PASS" if ok else "FAIL"), "-", label, "" if ok else ":: " + detail)
    return ok


# --------------------------------------------------------------------------- #
# No-network dummy OpenAI-compatible client (mirrors tests/test_v03_ai.py)
# --------------------------------------------------------------------------- #
class _MockResponse:
    def __init__(self, content):
        self._content = content

    class _Msg:
        def __init__(self, c):
            self.content = c

    class _Choice:
        def __init__(self, c):
            self.message = _MockResponse._Msg(c)

    @property
    def choices(self):
        return [_MockResponse._Choice(self._content)]


class MockOpenAIClient:
    class _Completions:
        def __init__(self, owner):
            self._owner = owner

        def create(self, model=None, messages=None, **kwargs):
            self._owner.calls.append((model, messages))
            if self._owner._exc is not None:
                raise self._owner._exc
            return _MockResponse(self._owner._text)

    class _Chat:
        def __init__(self, owner):
            self.completions = MockOpenAIClient._Completions(owner)

    def __init__(self, text="MOCK_TEXT", exc=None):
        self._text = text
        self._exc = exc
        self.calls = []
        self.chat = MockOpenAIClient._Chat(self)


def _fake_resp(code):
    class R:
        status_code = code
        headers = {}
        request = None
        is_success = False

    return R()


def _mk_openai_error(kind):
    if kind == "auth":
        return openai.AuthenticationError(
            message="bad key", response=_fake_resp(401), body=None
        )
    if kind == "ratelimit":
        return openai.RateLimitError(
            message="rate limited", response=_fake_resp(429), body=None
        )
    raise ValueError(kind)


# --------------------------------------------------------------------------- #
# Fixtures
# --------------------------------------------------------------------------- #
def _make_png(path):
    from PIL import Image

    Image.new("RGB", (64, 64), (255, 0, 0)).save(path, format="PNG")


def _make_scanned_pdf(path):
    import fitz
    from PIL import Image

    img = Image.new("RGB", (200, 200), (245, 245, 245))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    doc = fitz.open()
    page = doc.new_page(width=200, height=200)
    page.insert_image(fitz.Rect(0, 0, 200, 200), stream=buf.read())
    doc.save(str(path))
    doc.close()


def _make_docx(path, png):
    from docx import Document

    d = Document()
    d.add_paragraph("Hello Docx")
    d.add_picture(str(png))
    d.save(str(path))


def _make_pptx(path, png):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.add_picture(str(png), 0, 0)
    prs.save(str(path))


def _make_xlsx(path, png):
    from openpyxl import Workbook
    from openpyxl.drawing.image import Image as XLImage

    wb = Workbook()
    ws = wb.active
    ws["A1"] = "Hello Xlsx"
    ws.add_image(XLImage(str(png)), "B2")
    wb.save(str(path))


def main():
    ok = True
    tmp = Path(tempfile.mkdtemp(prefix="mddesk_v03_frozen_"))
    png = tmp / "img.png"
    scanned = tmp / "scanned.pdf"
    _make_png(png)
    _make_scanned_pdf(scanned)

    # 1. collection checks
    ok &= _check("IMPORT_OPENAI", openai is not None)
    ok &= _check("OCR_PLUGIN_AVAILABLE", is_ocr_plugin_available())

    # 3. AI-OFF equivalence vs v0.2 in the frozen binary.
    # Use a self-generated HTML fixture so the smoke has no dependency on an
    # external test_input.html (which is not bundled into the frozen EXE).
    fixture = tmp / "sample.html"
    fixture.write_text(
        "<html><body><h1>Title</h1><p>Hello frozen world.</p></body></html>",
        encoding="utf-8",
    )
    plain = MarkItDown().convert(str(fixture)).markdown
    disabled_md = MarkItDownFactory.create(EngineConfig.disabled()).convert(
        str(fixture)
    ).markdown
    ok &= _check(
        "AI_OFF_EQUALS_V02",
        plain == disabled_md and len(plain) > 0,
        f"len_plain={len(plain)} len_disabled={len(disabled_md)}",
    )

    # 4. frozen plugin path registers OCR converters
    cfg_ai = EngineConfig(ai_enabled=True, ai_endpoint="", ai_api_key="", ai_model="gpt-4o")
    cfg_ai._client = MockOpenAIClient(text="X")
    md_ai = MarkItDownFactory.create(cfg_ai)
    ocr_count = sum(
        1 for c in md_ai._converters if "OCR" in type(c.converter).__name__
    )
    ok &= _check("FROZEN_OCR_CONVERTERS_REGISTERED", ocr_count >= 4, f"count={ocr_count}")

    # 5. image description (no network)
    cfg_img = EngineConfig(ai_enabled=True, ai_endpoint="", ai_api_key="", ai_model="gpt-4o")
    cfg_img._client = MockOpenAIClient(text="MOCK_CAPTION")
    md_img = convert_file(str(png), engine_config=cfg_img)
    ok &= _check(
        "IMAGE_DESCRIPTION_FROZEN",
        "MOCK_CAPTION" in md_img,
        repr(md_img[:80]),
    )

    # 6. scanned-PDF OCR (no network)
    cfg_pdf = EngineConfig(ai_enabled=True, ai_endpoint="", ai_api_key="", ai_model="gpt-4o")
    cfg_pdf._client = MockOpenAIClient(text="OCR_TEXT")
    md_pdf = convert_file(str(scanned), engine_config=cfg_pdf)
    ok &= _check(
        "SCANNED_PDF_OCR_FROZEN",
        "OCR_TEXT" in md_pdf and "*[Image OCR]" in md_pdf,
        repr(md_pdf[:120]),
    )

    # 7. scanned-PDF OCR FAILURE -> stable *[OCR Error] marker in frozen bin
    docx_f = tmp / "doc.docx"
    pptx_f = tmp / "doc.pptx"
    xlsx_f = tmp / "doc.xlsx"
    _make_docx(docx_f, png)
    _make_pptx(pptx_f, png)
    _make_xlsx(xlsx_f, png)

    cfg_pdf_err = EngineConfig(
        ai_enabled=True, ai_endpoint="", ai_api_key="", ai_model="gpt-4o"
    )
    cfg_pdf_err._client = MockOpenAIClient(exc=_mk_openai_error("auth"))
    md_pdf_err = convert_file(str(scanned), engine_config=cfg_pdf_err)
    ok &= _check(
        "SCANNED_PDF_OCR_ERROR_MARKER_FROZEN",
        "*[OCR Error]" in md_pdf_err
        and "[End OCR Error]*" in md_pdf_err
        and "*[Image OCR]" not in md_pdf_err,
        repr(md_pdf_err[:160]),
    )

    # 8. DOCX / PPTX / XLSX OCR FAILURE -> stable *[OCR Error] marker (frozen)
    for label, fpath in (("DOCX", docx_f), ("PPTX", pptx_f), ("XLSX", xlsx_f)):
        cfg_oe = EngineConfig(
            ai_enabled=True, ai_endpoint="", ai_api_key="", ai_model="gpt-4o"
        )
        cfg_oe._client = MockOpenAIClient(exc=_mk_openai_error("auth"))
        md_oe = convert_file(str(fpath), engine_config=cfg_oe)
        ok &= _check(
            f"OCR_ERROR_MARKER_FROZEN_{label}",
            "*[OCR Error]" in md_oe and "[End OCR Error]*" in md_oe,
            repr(md_oe[:160]),
        )

    print()
    if ok:
        print("V03_FROZEN_PASS")
        sys.exit(0)
    else:
        print("V03_FROZEN_FAIL")
        sys.exit(1)


if __name__ == "__main__":
    main()
