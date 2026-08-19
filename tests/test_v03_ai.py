"""MdDesk v0.3 AI-enhanced conversion tests (Stage 6).

Covers the v0.3 feature set:
  * AI-OFF path is byte-for-byte equivalent to v0.2 (no plugins, no LLM)
  * OCR plugin discovery + EngineConfig.should_enable_ocr()
  * Built-in LLM image description via a shared OpenAI-compatible client
  * Scanned-PDF OCR + DOCX/PPTX/XLSX embedded-image OCR (vendored markitdown-ocr)
  * AI error classification: image-description failures surface as friendly
    Chinese messages (verifies the FileConversionException.attempts path)
  * OCR backend failures are embedded as classified markdown and do NOT crash

Run: python tests/test_v03_ai.py   (or: pytest tests/test_v03_ai.py)
"""

import io
import sys
import tempfile
from pathlib import Path
from unittest.mock import patch

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

from markitdown import MarkItDown
import openai

from src.engine_config import EngineConfig, is_ocr_plugin_available
from src.markitdown_factory import MarkItDownFactory
from src.converter import convert_file, map_exception, ConversionError
from src.settings import Settings
from src.file_entry import FileStatus

FIX = ROOT / "tests" / "fixtures"
TEST_INPUT = ROOT.parent / "test_input.html"


# --------------------------------------------------------------------------- #
# Fixtures                                                                     #
# --------------------------------------------------------------------------- #
def _make_png(path: Path, color=(255, 0, 0)) -> None:
    from PIL import Image

    Image.new("RGB", (64, 64), color).save(path, format="PNG")


def _make_scanned_pdf(path: Path) -> None:
    """A PDF that contains ONLY an image page and NO selectable text."""
    import fitz
    from PIL import Image

    img = Image.new("RGB", (200, 200), (245, 245, 245))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    doc = fitz.open()
    page = doc.new_page(width=200, height=200)
    page.insert_image(fitz.Rect(0, 0, 200, 200), stream=buf.read())
    doc.save(str(path))
    doc.close()


def _make_docx(path: Path, png: Path) -> None:
    from docx import Document

    d = Document()
    d.add_paragraph("Hello Docx")
    d.add_picture(str(png))
    d.save(str(path))


def _make_pptx(path: Path, png: Path) -> None:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.add_picture(str(png), 0, 0)
    prs.save(str(path))


def _make_xlsx(path: Path, png: Path) -> None:
    from openpyxl import Workbook
    from openpyxl.drawing.image import Image as XLImage

    wb = Workbook()
    ws = wb.active
    ws["A1"] = "Hello Xlsx"
    ws.add_image(XLImage(str(png)), "B2")
    wb.save(str(path))


# --------------------------------------------------------------------------- #
# Mock OpenAI-compatible client                                               #
# --------------------------------------------------------------------------- #
class _MockResponse:
    def __init__(self, content):
        self._content = content

    class _Msg:
        def __init__(self, c):
            self.content = c

    class _Choice:
        def __init__(self, c):
            self.message = _MockResponse._Msg(c)

    @property
    def choices(self):
        return [_MockResponse._Choice(self._content)]


class MockOpenAIClient:
    """Returns a canned caption/OCR text, or raises a chosen exception."""

    class _Completions:
        def __init__(self, owner):
            self._owner = owner

        def create(self, model=None, messages=None, **kwargs):
            self._owner.calls.append((model, messages))
            if self._owner._exc is not None:
                raise self._owner._exc
            return _MockResponse(self._owner._text)

    class _Chat:
        def __init__(self, owner):
            self.completions = MockOpenAIClient._Completions(owner)

    def __init__(self, text="MOCK_TEXT", exc=None):
        self._text = text
        self._exc = exc
        self.calls = []
        self.chat = MockOpenAIClient._Chat(self)


def _mk_openai_error(kind: str):
    """Build a real openai SDK exception (no network) for `kind`."""
    if kind == "auth":
        return openai.AuthenticationError(
            message="bad key", response=_fake_resp(401), body=None
        )
    if kind == "permission":
        return openai.PermissionDeniedError(
            message="no permission", response=_fake_resp(403), body=None
        )
    if kind == "ratelimit":
        return openai.RateLimitError(
            message="rate limited", response=_fake_resp(429), body=None
        )
    if kind == "badrequest":
        return openai.BadRequestError(
            message="bad request", response=_fake_resp(400), body=None
        )
    if kind == "connection":
        return openai.APIConnectionError(message="conn fail", request=None)
    if kind == "timeout":
        return openai.APITimeoutError(request=None)
    raise ValueError(kind)


def _fake_resp(code: int):
    class R:
        status_code = code
        headers = {}
        request = None
        is_success = False

    return R()


# --------------------------------------------------------------------------- #
# Check helper                                                                #
# --------------------------------------------------------------------------- #
def _check(name, cond, detail=""):
    print(("PASS" if cond else "FAIL"), "-", name, "" if cond else f":: {detail}")
    return cond


# --------------------------------------------------------------------------- #
# Test body                                                                   #
# --------------------------------------------------------------------------- #
def _run_all() -> bool:
    ok = True
    tmp = Path(tempfile.mkdtemp(prefix="mddesk_v03_"))
    png = tmp / "img.png"
    scanned = tmp / "scanned.pdf"
    docx = tmp / "doc.docx"
    pptx = tmp / "doc.pptx"
    xlsx = tmp / "doc.xlsx"
    _make_png(png)
    _make_scanned_pdf(scanned)
    _make_docx(docx, png)
    _make_pptx(pptx, png)
    _make_xlsx(xlsx, png)

    # 1. AI-OFF equivalence vs v0.2 (regression) -------------------------------
    md_disabled = MarkItDownFactory.create(EngineConfig.disabled())
    ok &= _check(
        "1. AI-OFF 返回纯 MarkItDown 实例",
        isinstance(md_disabled, MarkItDown),
        type(md_disabled).__name__,
    )
    has_ocr = any("OCR" in type(c.converter).__name__ for c in md_disabled._converters)
    ok &= _check("1b. AI-OFF 未注册 OCR 转换器", not has_ocr)

    plain = MarkItDown().convert(str(TEST_INPUT)).markdown
    disabled_md = MarkItDownFactory.create(EngineConfig.disabled()).convert(
        str(TEST_INPUT)
    ).markdown
    ok &= _check(
        "2. AI-OFF 与 v0.2 输出逐字一致",
        plain == disabled_md,
        f"len plain={len(plain)} disabled={len(disabled_md)}",
    )

    # 3. OCR plugin discovery + EngineConfig ----------------------------------
    ok &= _check("3. markitdown-ocr 插件可发现", is_ocr_plugin_available())
    s_on = Settings.default()
    s_on.ai.enabled = True
    s_on.ai.model = "gpt-4o"
    cfg_on = EngineConfig.from_settings(s_on)
    ok &= _check("3b. AI 启用时 ai_enabled", cfg_on.ai_enabled)
    ok &= _check("3c. 插件可用 -> should_enable_ocr", cfg_on.should_enable_ocr())
    cfg_off = EngineConfig.from_settings(Settings.default())
    ok &= _check("3d. AI 关闭 -> 不启用 OCR", not cfg_off.should_enable_ocr())

    # 4. Image description via shared client (mock) ---------------------------
    cfg_img = EngineConfig(
        ai_enabled=True, ai_endpoint="", ai_api_key="", ai_model="gpt-4o"
    )
    cfg_img._client = MockOpenAIClient(text="MOCK_CAPTION")
    md_img = convert_file(str(png), engine_config=cfg_img)
    ok &= _check(
        "4. 图片描述注入 Markdown",
        "MOCK_CAPTION" in md_img,
        repr(md_img[:80]),
    )

    # 5. Scanned-PDF OCR via mock ---------------------------------------------
    cfg_pdf = EngineConfig(
        ai_enabled=True, ai_endpoint="", ai_api_key="", ai_model="gpt-4o"
    )
    client_pdf = MockOpenAIClient(text="OCR_TEXT")
    cfg_pdf._client = client_pdf
    md_pdf = convert_file(str(scanned), engine_config=cfg_pdf)
    ok &= _check(
        "5. 扫描件 PDF OCR 文本出现",
        "OCR_TEXT" in md_pdf and "*[Image OCR]" in md_pdf,
        repr(md_pdf[:120]),
    )
    ok &= _check("5b. OCR 后端确实被调用", len(client_pdf.calls) > 0,
                 f"calls={len(client_pdf.calls)}")

    # 6. AI error classification + v0.6 error isolation --
    # v0.6 contract change (plan §3.4): an AI failure during image
    # description must NO LONGER fail the whole conversion. The converter
    # downgrades to the non-AI path and attaches an AI_PROVIDER_FAILURE
    # warning carrying the classified Chinese message.
    for kind, needle in (
        ("auth", "鉴权失败"),
        ("ratelimit", "额度或频率"),
        ("connection", "连接失败"),
        ("badrequest", "请求被拒绝"),
    ):
        cfg_e = EngineConfig(
            ai_enabled=True, ai_endpoint="", ai_api_key="", ai_model="gpt-4o"
        )
        cfg_e._client = MockOpenAIClient(exc=_mk_openai_error(kind))
        warns6 = []
        md6 = convert_file(str(png), engine_config=cfg_e, warnings_out=warns6)
        ok &= _check(
            f"6.{kind} 降级成功(不抛错)",
            isinstance(md6, str),
            "unexpected ConversionError",
        )
        ok &= _check(
            f"6.{kind} 友好提示({needle})",
            any(
                w.code == "AI_PROVIDER_FAILURE" and needle in w.message
                for w in warns6
            ),
            f"warnings={[w.code for w in warns6]}",
        )
        ok &= _check(
            f"6.{kind} 提示不含模型输出描述",
            "Description:" not in md6,
            "downgraded output still carries a description",
        )

    # 7. OCR backend failure -> embedded classified markdown, no crash --------
    cfg_ocr_err = EngineConfig(
        ai_enabled=True, ai_endpoint="", ai_api_key="", ai_model="gpt-4o"
    )
    cfg_ocr_err._client = MockOpenAIClient(exc=_mk_openai_error("auth"))
    md_ocr_err = convert_file(str(scanned), engine_config=cfg_ocr_err)
    ok &= _check(
        "7. OCR 失败不崩溃(返回 Markdown)",
        isinstance(md_ocr_err, str) and len(md_ocr_err) > 0,
    )
    ok &= _check(
        "7b. OCR 失败内联分类提示",
        "OCR 失败：API Key 无效" in md_ocr_err,
        repr(md_ocr_err[:160]),
    )
    ok &= _check(
        "7c. OCR 错误块标记稳定(*[OCR Error])",
        "*[OCR Error]" in md_ocr_err and "[End OCR Error]*" in md_ocr_err,
        repr(md_ocr_err[:160]),
    )
    ok &= _check(
        "7d. OCR 错误块未被伪装成成功块",
        "*[Image OCR]" not in md_ocr_err,
        repr(md_ocr_err[:160]),
    )

    # 8. DOCX / PPTX / XLSX embedded-image OCR --------------------------------
    for label, fpath, token in (
        ("DOCX", docx, "DOCX_OCR"),
        ("PPTX", pptx, "PPTX_OCR"),
        ("XLSX", xlsx, "XLSX_OCR"),
    ):
        cfg_o = EngineConfig(
            ai_enabled=True, ai_endpoint="", ai_api_key="", ai_model="gpt-4o"
        )
        client_o = MockOpenAIClient(text=token)
        cfg_o._client = client_o
        md_o = convert_file(str(fpath), engine_config=cfg_o)
        ok &= _check(
            f"8.{label} 内嵌图 OCR 文本出现",
            token in md_o and "*[Image OCR]" in md_o,
            repr(md_o[:120]),
        )
        ok &= _check(
            f"8b.{label} OCR 后端被调用",
            len(client_o.calls) > 0,
            f"calls={len(client_o.calls)}",
        )

    # 9. DOCX / PPTX / XLSX OCR FAILURE -> stable *[OCR Error] block (NOT
    #    dropped, NOT disguised as success) --------------------------------
    for label, fpath in (
        ("DOCX", docx),
        ("PPTX", pptx),
        ("XLSX", xlsx),
    ):
        cfg_oe = EngineConfig(
            ai_enabled=True, ai_endpoint="", ai_api_key="", ai_model="gpt-4o"
        )
        cfg_oe._client = MockOpenAIClient(exc=_mk_openai_error("auth"))
        md_oe = convert_file(str(fpath), engine_config=cfg_oe)
        ok &= _check(
            f"9.{label} OCR 失败内联错误块",
            "*[OCR Error]" in md_oe and "[End OCR Error]*" in md_oe,
            repr(md_oe[:160]),
        )
        ok &= _check(
            f"9b.{label} OCR 错误块未伪装成成功",
            "*[Image OCR]" not in md_oe,
            repr(md_oe[:160]),
        )

    return ok


def test_v03_ai_suite():
    assert _run_all()


def main():
    ok = _run_all()
    print()
    print("ALL v0.3 AI CHECKS PASSED" if ok else "SOME v0.3 AI CHECKS FAILED")
    sys.exit(0 if ok else 1)


if __name__ == "__main__":
    main()
