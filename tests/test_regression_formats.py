#!/usr/bin/env python
"""Stage 1B regression across previously-PASSING formats.

Generates a tiny valid sample for each format that Stage 1A confirmed as
PASS, converts it through the MarkItDown engine, and asserts a non-empty
Markdown result (no crash, no regression introduced by the MSG/olefile
change). This is pure verification — no converter code is modified.

Run from project root:
    python tests/test_regression_formats.py
"""

import csv
import io
import json
import os
import sys
import tempfile
import zipfile
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

TEST_INPUT_HTML = ROOT.parent / "test_input.html"


def _check(name, cond, detail=""):
    print(("PASS" if cond else "FAIL"), "-", name, "" if cond else f":: {detail}")
    return cond


def _write(path: Path, data) -> str:
    if isinstance(data, str):
        data = data.encode("utf-8")
    path.write_bytes(data)
    return str(path)


def build_samples(d: Path):
    samples = {}

    # TXT
    samples["TXT"] = _write(d / "s.txt", "# Title\n\nhello *world*")

    # JSON
    samples["JSON"] = _write(d / "s.json", json.dumps({"a": 1, "b": [2, 3]}, ensure_ascii=False))

    # CSV
    buf = io.StringIO()
    w = csv.writer(buf)
    w.writerow(["name", "age"])
    w.writerow(["Alice", "30"])
    samples["CSV"] = _write(d / "s.csv", buf.getvalue())

    # XML (generic)
    samples["XML"] = _write(d / "s.xml", '<?xml version="1.0"?><root><item>value</item></root>')

    # HTML
    if TEST_INPUT_HTML.exists():
        samples["HTML"] = str(TEST_INPUT_HTML)

    # EPUB (zip with mimetype + OPF + xhtml)
    ep = d / "s.epub"
    with zipfile.ZipFile(ep, "w") as z:
        z.writestr("mimetype", "application/epub+zip")
        z.writestr("META-INF/container.xml",
                   '<?xml version="1.0"?><container version="1.0" '
                   'xmlns="urn:oasis:names:tc:opendocument:xmlns:container">'
                   '<rootfiles><rootfile full-path="OEBPS/content.opf" '
                   'media-type="application/oebps-package+xml"/></rootfiles>'
                   '</container>')
        z.writestr("OEBPS/content.opf",
                   '<?xml version="1.0"?><package xmlns="http://www.idpf.org/2007/opf" '
                   'version="2.0" unique-identifier="id"><metadata '
                   'xmlns:dc="http://purl.org/dc/elements/1.1/">'
                   '<dc:title>Epub Title</dc:title><dc:creator>Epub Author</dc:creator>'
                   '<dc:language>en</dc:language></metadata>'
                   '<manifest><item id="c1" href="content.xhtml" '
                   'media-type="application/xhtml+xml"/></manifest>'
                   '<spine><itemref idref="c1"/></spine></package>')
        z.writestr("OEBPS/content.xhtml",
                   "<html><body><h1>Epub Heading</h1><p>Epub body text.</p></body></html>")
    samples["EPUB"] = str(ep)

    # ZIP (containing a txt)
    zp = d / "s.zip"
    with zipfile.ZipFile(zp, "w") as z:
        z.writestr("inside.txt", "zip inner content")
    samples["ZIP"] = str(zp)

    # PDF
    try:
        from fpdf import FPDF
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Helvetica", size=20)
        pdf.cell(0, 10, "PDF Title")
        pdf.output(str(d / "s.pdf"))
        samples["PDF"] = str(d / "s.pdf")
    except Exception as exc:  # pragma: no cover
        print("SKIP - PDF fixture:", exc)

    # DOCX
    try:
        from docx import Document
        doc = Document()
        doc.add_heading("Docx Heading", 0)
        doc.add_paragraph("docx body")
        doc.save(str(d / "s.docx"))
        samples["DOCX"] = str(d / "s.docx")
    except Exception as exc:  # pragma: no cover
        print("SKIP - DOCX fixture:", exc)

    # XLSX
    try:
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws["A1"] = "X"
        ws["B1"] = "Y"
        ws.append([1, 2])
        wb.save(str(d / "s.xlsx"))
        samples["XLSX"] = str(d / "s.xlsx")
    except Exception as exc:  # pragma: no cover
        print("SKIP - XLSX fixture:", exc)

    # PPTX
    try:
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Pptx Title"
        prs.save(str(d / "s.pptx"))
        samples["PPTX"] = str(d / "s.pptx")
    except Exception as exc:  # pragma: no cover
        print("SKIP - PPTX fixture:", exc)

    # IPYNB
    nb = {
        "cells": [{"cell_type": "markdown", "metadata": {}, "source": ["# Notebook\n", "markdown cell"]},
                  {"cell_type": "code", "execution_count": None, "metadata": {}, "outputs": [], "source": ["print(1)"]}],
        "metadata": {"kernelspec": {"display_name": "Python", "language": "python", "name": "python"}},
        "nbformat": 4, "nbformat_minor": 5,
    }
    samples["IPYNB"] = _write(d / "s.ipynb", json.dumps(nb, ensure_ascii=False))

    # Image (PNG)
    try:
        from PIL import Image
        im = Image.new("RGB", (32, 32), (255, 0, 0))
        im.save(str(d / "s.png"))
        samples["PNG"] = str(d / "s.png")
    except Exception as exc:  # pragma: no cover
        print("SKIP - PNG fixture:", exc)

    return samples


def main():
    from markitdown import MarkItDown
    d = Path(tempfile.mkdtemp(prefix="reg_"))
    samples = build_samples(d)
    md = MarkItDown()
    ok = True
    # Formats that yield text by default. PNG is exempt: MarkItDown produces no
    # text for images unless the OCR extra is installed (unchanged behavior).
    hard = ["TXT", "JSON", "CSV", "XML", "HTML", "EPUB", "ZIP", "PDF", "DOCX", "XLSX", "PPTX", "IPYNB"]
    exempt = {"PNG"}
    for fmt in hard + list(exempt):
        if fmt not in samples:
            print("SKIP -", fmt, "(fixture unavailable)")
            continue
        try:
            out = md.convert(samples[fmt]).markdown
        except Exception as exc:  # noqa: BLE001
            ok &= _check(f"REGRESSION {fmt}", False, f"exception {exc!r}")
            continue
        non_empty = bool(out and out.strip())
        if fmt in exempt:
            # No-text is the expected, unchanged behavior for images.
            _check(f"REGRESSION {fmt} (no text without OCR extra; unchanged)", True,
                   f"len={len(out) if out else 0}")
            print(f"    [{fmt}] md_len={len(out) if out else 0} (expected empty; OCR extra absent)")
            continue
        ok &= _check(f"REGRESSION {fmt}", non_empty, f"len={len(out) if out else 0}")
        if non_empty:
            print(f"    [{fmt}] md_len={len(out)} preview={out.strip()[:48]!r}")
    print()
    print("ALL FORMAT REGRESSION CHECKS PASSED" if ok else "SOME FORMAT REGRESSION CHECKS FAILED")
    sys.exit(0 if ok else 1)


if __name__ == "__main__":
    main()
